from __future__ import annotations

import json
import csv
import hashlib
import os
import re
import shutil
import subprocess
import tempfile
import time
import zipfile
from pathlib import Path
from typing import Iterable

from .deliverables import (
    ODF_ARTIFACT_EXTENSIONS,
    OLE_ARTIFACT_EXTENSIONS,
    OOXML_ARTIFACT_MARKERS,
    PROFESSIONAL_ARTIFACT_EXTENSIONS,
    SUPPORT_ARTIFACT_DIR_NAMES,
    candidate_artifact_paths,
    is_user_deliverable_relative_path,
    is_valid_professional_artifact,
)
from .failure_classification import classify_cli_failure, has_structured_failure_evidence
from .runtime_identity import derive_legacy_backend_label


RUN_EVIDENCE_DIR_NAME = "glasshive-run"
_MAX_TEXT_SCAN_BYTES = 512_000
_MAX_TEXT_SCAN_FILES = 200
_SECRET_KEY_MARKERS = (
    "SECRET",
    "TOKEN",
    "KEY",
    "COOKIE",
    "PASSWORD",
    "PASSWD",
    "PWD",
    "CREDENTIAL",
    "AUTH",
)
_SECRET_REDACTIONS: tuple[tuple[re.Pattern[str], str], ...] = (
    (re.compile(r"/Users/[^/\s\"'`]+(?:/[^\s\"'`]+)+"), "[REDACTED_LOCAL_PATH]"),
    (re.compile(r"~/[^\s\"'`]+(?:/[^\s\"'`]+)+"), "[REDACTED_LOCAL_PATH]"),
    (re.compile(r"/Users/[^/\s\"']+"), "~"),
    (re.compile(r"(?i)(bearer\s+)[A-Za-z0-9._~+/=-]{12,}"), r"\1[REDACTED]"),
    (re.compile(r"(?i)((?:api[_-]?key|token|secret|password|passwd|pwd)\s*[:=]\s*)[^\s\"']{6,}"), r"\1[REDACTED]"),
    (re.compile(r"(?i)([?&](?:gh_token|gh_sig|gh_exp|gh_kind)=)([^&#\s\"']+)"), r"\1[REDACTED]"),
    (re.compile(r"\bsk-[A-Za-z0-9_-]{12,}\b"), "sk-[REDACTED]"),
    (re.compile(r"\b[A-Za-z0-9_]{8,}:[A-Za-z0-9_./+=-]{20,}\b"), "[REDACTED_CREDENTIAL]"),
)
_FINAL_REPORT_RE = re.compile(
    r"(?m)^[ \t]*(?:#{1,6}[ \t]+|>[ \t]*)?(?:[*_]{1,3})?FINAL REPORT\s*:\s*(?:[*_]{1,3})?",
    re.I,
)

_MONTHS = {
    "jan": 1,
    "january": 1,
    "feb": 2,
    "february": 2,
    "mar": 3,
    "march": 3,
    "apr": 4,
    "april": 4,
    "may": 5,
    "jun": 6,
    "june": 6,
    "jul": 7,
    "july": 7,
    "aug": 8,
    "august": 8,
    "sep": 9,
    "sept": 9,
    "september": 9,
    "oct": 10,
    "october": 10,
    "nov": 11,
    "november": 11,
    "dec": 12,
    "december": 12,
}
_MONTH_PATTERN = (
    "jan(?:uary)?|feb(?:ruary)?|mar(?:ch)?|apr(?:il)?|may|jun(?:e)?|jul(?:y)?|"
    "aug(?:ust)?|sep(?:t(?:ember)?|tember)?|oct(?:ober)?|nov(?:ember)?|dec(?:ember)?"
)
_MONTH_YEAR_RE = re.compile(rf"\b({_MONTH_PATTERN})\s+(\d{{4}})\b", re.IGNORECASE)
_ISO_YEAR_MONTH_RE = re.compile(r"\b(20\d{2})-(0[1-9]|1[0-2])\b")
_DATE_LIMIT_RE = re.compile(
    rf"\b(?:through|until|no later than|ending|ended|before end of|up to|as of|by)\s+({_MONTH_PATTERN})\s+(\d{{4}})\b",
    re.IGNORECASE,
)

_CONTENT_HYGIENE_PATTERNS: tuple[tuple[re.Pattern[str], str], ...] = (
    (re.compile(r"skip\s+(?:to\s+)?(?:main\s+)?(?:content|navigation)", re.I), "navigation chrome"),
    (re.compile(r"cookie\s+(?:settings|preferences|policy|notice|banner)", re.I), "cookie banner"),
    (re.compile(r"privacy\s+policy|terms\s+(?:of\s+(?:use|service)|and\s+conditions)", re.I), "legal chrome"),
    (re.compile(r"all\s+rights\s+reserved|(?:read|view)\s+(?:more|all)", re.I), "page chrome"),
    (re.compile(r"please\s+enable\s+(?:java\s*)?script|enable\s+js", re.I), "javascript wall"),
    (re.compile(r"subscribe\s+to\s+continue|sign\s+in\s+to\s+continue|paywall", re.I), "paywall/auth wall"),
    (re.compile(r"are\s+you\s+(?:a\s+)?robot|captcha|bot\s+wall|403\s+forbidden|access\s+denied", re.I), "bot/access wall"),
    (re.compile(r"sourceurl=|__nuxt__|@layer|tailwindcss", re.I), "raw web asset"),
    (re.compile(r"\b(?:window|document)\.[A-Za-z_$]"), "javascript"),
    (re.compile(r"\bfunction\s+(?:[A-Za-z_$][\w$]*)?\s*\([^)]*\)\s*\{", re.I), "javascript"),
    (re.compile(r"\bvar\s+[A-Za-z_$][\w$]*\s*="), "javascript"),
    (re.compile(r"&(?:nbsp|amp|lt|gt|quot|apos);|&#\d+;|&#x[0-9a-f]+;", re.I), "html entity"),
)
_CONTENT_HYGIENE_NEGATIVE_CONTEXT_RE = re.compile(
    r"\b(?:exclude|excludes|excluded|free\s+of|without|no|not)\b[^.\n;]{0,120}$",
    re.I,
)
_TEXT_FILE_SUFFIXES = {".csv", ".json", ".jsonl", ".md", ".txt", ".html", ".htm", ".tsv"}
_CONTENT_HYGIENE_TEXT_SUFFIXES = {".csv", ".json", ".jsonl", ".md", ".txt", ".tsv"}
_HTML_FILE_SUFFIXES = {".html", ".htm"}
_CONSTRAINT_BINARY_TEXT_SUFFIXES = {".pdf", ".xlsx", ".docx", ".pptx"}
_OUTPUT_FORMAT_SUFFIXES = (
    "pdf",
    "xlsx",
    "xls",
    "csv",
    "html",
    "docx",
    "pptx",
    "txt",
    "md",
    "json",
    "tsv",
    "rtf",
    "odt",
    "ods",
    "odp",
)
_OUTPUT_FORMAT_ALIASES = {
    "xlsx": {"xlsx", "xls", "ods"},
    "xls": {"xlsx", "xls", "ods"},
    "csv": {"csv", "tsv"},
    "html": {"html", "htm"},
    "docx": {"docx", "doc", "odt", "rtf"},
    "pptx": {"pptx", "ppt", "odp"},
    "txt": {"txt"},
    "md": {"md"},
}
_OUTPUT_ACTION_RE = re.compile(
    r"\b(deliver|create|generate|produce|write|save|export|output|return|build|prepare|compose|draft|render|"
    r"artifact|artifacts|deliverable|deliverables|report|reports|file|files)\b",
    re.I,
)
_FINAL_ANSWER_ONLY_CONTEXT_RE = re.compile(
    r"\b(in\s+chat|inline|in\s+(?:the\s+)?final\s+answer|answer\s+(?:here|inline|in\s+chat)|"
    r"respond\s+(?:here|inline|in\s+chat)|reply\s+(?:here|inline|in\s+chat)|"
    r"tell\s+me|report\s+back)\b",
    re.I,
)
_FILE_DELIVERABLE_CONTEXT_RE = re.compile(
    r"\b(file|files|artifact|artifacts|deliverable|deliverables|save|export|download|"
    r"pdf|xlsx|xls|csv|html|docx|pptx|json|tsv|workbook|spreadsheet|slide\s+deck|deck)\b",
    re.I,
)
_OUTPUT_FORBIDDEN_RE = re.compile(
    r"\b(do\s+not|don't|must\s+not|should\s+not|without|avoid|forbidden|no)\b"
    r"[^.\n;]{0,80}\b(create|generate|produce|write|save|export|output|deliver|artifact|file|report|pdf|xlsx|xls|csv|html|docx|pptx|txt|md|json|tsv)\b",
    re.I,
)
_OUTPUT_FORBIDDEN_CLAUSE_RE = re.compile(
    r"\b(?:do\s+not|don't|must\s+not|should\s+not|without|avoid|forbidden|no)\b[^.\n;)]*",
    re.I,
)
_COVERAGE_ENTITY_RE = r"(?:firms?|companies|entities|targets?|items|records|rows|entries)"
_COVERAGE_RANGE_RE = re.compile(
    rf"\b(?:aim\s+for|target|cover|include|research|score|screen|produce)\b[^.\n;]{{0,100}}"
    rf"\b(\d{{1,4}})\s*(?:-|–|—|to)\s*(\d{{1,4}})\s+({_COVERAGE_ENTITY_RE})\b",
    re.I,
)
_COVERAGE_RANGE_TOTAL_RE = re.compile(
    rf"\b(\d{{1,4}})\s*(?:-|–|—|to)\s*(\d{{1,4}})\s+({_COVERAGE_ENTITY_RE})\b"
    rf"[^.\n;]{{0,80}}\b(?:total|before\s+scoring|scored|screened)\b",
    re.I,
)
_COVERAGE_RANGE_BARE_TOTAL_RE = re.compile(
    r"\b(?:aim\s+for|target|cover|include|research|score|screen|produce)\b[^.\n;]{0,80}"
    r"\b(\d{1,4})\s*(?:-|–|—|to)\s*(\d{1,4})\s+(?:total|before\s+scoring|scored|screened)\b",
    re.I,
)
_COVERAGE_MIN_RE = re.compile(
    rf"\b(?:at\s+least|minimum(?:\s+of)?|min(?:imum)?|no\s+fewer\s+than)\s+(\d{{1,4}})\s+({_COVERAGE_ENTITY_RE})\b",
    re.I,
)
_COVERAGE_PLUS_RE = re.compile(
    rf"\b(\d{{1,4}})\+\s+({_COVERAGE_ENTITY_RE})\b",
    re.I,
)
_INPUT_FORMAT_CONTEXT_RE = re.compile(
    r"\b(attached|uploaded|input|source|provided|existing|original|reference|marker|embedded|"
    r"bytes?|signature|header|first[-\s]*byte|workspace|read|compare|analy[sz]e|from|using|use)\b",
    re.I,
)
_SUPPORT_ARTIFACT_DIRS = SUPPORT_ARTIFACT_DIR_NAMES
_RAW_SUPPORT_CAPTURE_DIRS = {
    "crawl_cache",
    "raw_pages",
    "raw_web",
    "site_snapshots",
    "source_snapshots",
    "web_snapshots",
}
_SEED_DESCRIPTOR_WORDS = {
    "attached",
    "companies",
    "entities",
    "examples",
    "firm",
    "files",
    "firms",
    "input",
    "inputs",
    "items",
    "list",
    "provided",
    "source",
    "targets",
    "topic",
    "topics",
    "uploaded",
}
_SKIP_SCAN_DIRS = {
    ".git",
    ".venv",
    "__pycache__",
    "chrome-user-data",
    "chromium-user-data",
    RUN_EVIDENCE_DIR_NAME,
    "node_modules",
}


def _now_epoch() -> float:
    return time.time()


def _redact_text(value: object, *, max_chars: int | None = None) -> str:
    text = str(value or "")
    for pattern, replacement in _SECRET_REDACTIONS:
        text = pattern.sub(replacement, text)
    if max_chars is not None and len(text) > max_chars:
        return text[-max_chars:]
    return text


def _redact_command_arg(value: object) -> str:
    text = str(value or "")
    if "\n" in text or len(text) > 600:
        digest = hashlib.sha256(text.encode("utf-8", errors="replace")).hexdigest()[:16]
        return f"[REDACTED_LONG_ARG chars={len(text)} sha256={digest}]"
    if text.startswith("/") or text.startswith("~/"):
        return Path(text).name or "[REDACTED_PATH]"
    return _redact_text(text)


def _stdout_agent_has_final_report(stdout_text: str) -> bool:
    final_text_keys = {
        "finalAssistantVisibleText",
        "final_assistant_visible_text",
        "finalVisibleText",
        "final_visible_text",
    }

    def iter_dicts(value: object) -> Iterable[dict[str, object]]:
        if isinstance(value, dict):
            yield value
            for nested in value.values():
                yield from iter_dicts(nested)
        elif isinstance(value, list):
            for nested in value:
                yield from iter_dicts(nested)

    for line in str(stdout_text or "").splitlines():
        try:
            payload = json.loads(line)
        except json.JSONDecodeError:
            continue
        if not isinstance(payload, dict):
            continue
        for item in iter_dicts(payload):
            for key in final_text_keys:
                text = str(item.get(key) or "")
                if _FINAL_REPORT_RE.search(text):
                    return True
        if str(payload.get("type") or "") == "result":
            result_text = str(payload.get("result") or "")
            if _FINAL_REPORT_RE.search(result_text):
                return True
        item = payload.get("item")
        if not isinstance(item, dict):
            continue
        if str(item.get("type") or "") not in {"agent_message", "assistant_message"}:
            continue
        text = str(item.get("text") or "")
        if _FINAL_REPORT_RE.search(text):
            return True
    return bool(_FINAL_REPORT_RE.search(str(stdout_text or "")))


def _safe_env_keys(env: dict[str, str] | None) -> list[str]:
    keys: list[str] = []
    for key in sorted((env or {}).keys()):
        upper = key.upper()
        if any(marker in upper for marker in _SECRET_KEY_MARKERS):
            continue
        keys.append(key)
    return keys


def _effective_effort(worker: dict[str, object], command: list[str], env: dict[str, str] | None) -> str:
    configured = str(worker.get("effort") or worker.get("reasoning_effort") or "").strip()
    if configured:
        return configured
    for key in ("WPR_CODEX_CLI_REASONING_EFFORT", "CODEX_REASONING_EFFORT", "WPR_CLAUDE_CODE_EFFORT"):
        value = str((env or {}).get(key) or "").strip()
        if value:
            return value
    for index, part in enumerate(command):
        text = str(part or "").strip()
        if text.startswith("model_reasoning_effort"):
            return text.split("=", 1)[1].strip().strip("\"'")
        if text == "-c" and index + 1 < len(command):
            candidate = str(command[index + 1] or "").strip()
            if candidate.startswith("model_reasoning_effort"):
                return candidate.split("=", 1)[1].strip().strip("\"'")
        if text == "--effort" and index + 1 < len(command):
            return str(command[index + 1] or "").strip()
    return ""


def _transcript_metadata(
    transcript_paths: dict[str, str] | None,
    *,
    workspace_dir: Path | str | None = None,
    tail_limit_chars: int = 4000,
) -> dict[str, object]:
    workspace = Path(workspace_dir) if workspace_dir is not None else None
    metadata: dict[str, object] = {}
    for key, raw_path in (transcript_paths or {}).items():
        path_text = str(raw_path or "").strip()
        if not path_text:
            metadata[str(key)] = {"exists": False, "bytes": 0, "tail_limit_chars": tail_limit_chars, "truncated": False}
            continue
        path = Path(path_text)
        if workspace is not None and not path.is_absolute():
            path = workspace / path
        try:
            size = path.stat().st_size if path.is_file() else 0
            exists = path.is_file()
        except OSError:
            size = 0
            exists = False
        metadata[str(key)] = {
            "exists": exists,
            "bytes": size,
            "tail_limit_chars": tail_limit_chars if str(key) in {"stdout", "stderr"} else None,
            "truncated": bool(exists and str(key) in {"stdout", "stderr"} and size > tail_limit_chars),
            "capture_source": "file" if exists else "missing",
        }
    return metadata


def _effort_projection(worker: dict[str, object], effective_effort: str) -> dict[str, object]:
    raw = worker.get("effort_projection")
    if not isinstance(raw, dict):
        raw = worker.get("_effort_projection")
    if isinstance(raw, dict):
        allowed = raw.get("allowed")
        return {
            "requested": str(raw.get("requested") or ""),
            "effective": str(raw.get("effective") or effective_effort or ""),
            "allowed": [str(item) for item in allowed] if isinstance(allowed, list) else [],
            "route_proven": bool(raw.get("route_proven")),
            "fallback_reason": str(raw.get("fallback_reason") or ""),
        }
    return {
        "requested": effective_effort,
        "effective": effective_effort,
        "allowed": [],
        "route_proven": False,
        "fallback_reason": "",
    }


def _derived_legacy_backend(worker: dict[str, object]) -> str:
    return derive_legacy_backend_label(
        profile=worker.get("profile"),
        runtime=worker.get("runtime"),
        backend=worker.get("backend"),
    )


def _lines(text: str) -> list[str]:
    return [line.strip(" \t-*") for line in str(text or "").splitlines() if line.strip(" \t-*")]


def _contains_any(text: str, needles: tuple[str, ...]) -> bool:
    lower = text.lower()
    return any(needle in lower for needle in needles)


_DATE_WORD_RE = re.compile(r"\b(?:date|dated)\b", re.I)
_DATE_RANGE_WORD_RE = re.compile(r"\b(?:through|until|no later than|from|between)\b", re.I)
_YEAR_RE = re.compile(r"\b20\d{2}\b")
_SOURCE_DATE_PLANNING_CONTEXT_RE = re.compile(
    r"\b(?:sources|citation|citations|cited|primary|official|publication|published|dataset)\b",
    re.I,
)
_EVIDENCE_SOURCE_CONTEXT_RE = re.compile(
    r"\b(?:public|web|official|primary|external|published|cited)\s+evidence\b|\bevidence\s+(?:from|in|published|dated|cited)\b",
    re.I,
)


def _line_has_date_constraint_context(line: str) -> bool:
    text = str(line or "")
    has_date_value = bool(_MONTH_YEAR_RE.search(text) or _ISO_YEAR_MONTH_RE.search(text) or _YEAR_RE.search(text))
    return bool(_DATE_WORD_RE.search(text) or (_DATE_RANGE_WORD_RE.search(text) and has_date_value))


def _line_requires_source_date_planning_preservation(line: str) -> bool:
    text = str(line or "")
    return bool(
        _line_has_date_constraint_context(text)
        or _SOURCE_DATE_PLANNING_CONTEXT_RE.search(text)
        or _EVIDENCE_SOURCE_CONTEXT_RE.search(text)
    )


def _memory_write_mode_is_off(text: str) -> bool:
    return bool(re.search(r"\bmemory\s+write\s+mode\s*:\s*off\b", str(text or ""), re.I))


def _line_is_memory_proposal_output_context(line: str) -> bool:
    return bool(re.search(r"\bmemory[-\s]+proposals?\b", str(line or ""), re.I))


def _is_internal_scheduled_prompt_snapshot_path(path: str) -> bool:
    try:
        parts = Path(str(path or "")).parts
    except TypeError:
        return False
    return bool(parts and parts[0].lower() == "scheduled-prompt")


def _bootstrap_seed_files(worker: dict[str, object]) -> list[str]:
    try:
        bundle = json.loads(str(worker.get("bootstrap_bundle_json") or "{}"))
    except json.JSONDecodeError:
        bundle = {}
    if not isinstance(bundle, dict):
        bundle = {}
    files: list[str] = []
    for key in ("files", "input_files", "uploads"):
        raw = bundle.get(key)
        if not isinstance(raw, list):
            continue
        for item in raw:
            if isinstance(item, dict):
                candidate = str(item.get("path") or item.get("name") or item.get("relative_path") or "").strip()
            else:
                candidate = str(item or "").strip()
            if candidate and not _is_internal_scheduled_prompt_snapshot_path(candidate):
                files.append(f"uploaded file {candidate}")
    return files


def _is_markdown_heading(line: str) -> bool:
    text = str(line or "").strip()
    return bool(re.match(r"^#{1,6}\s+", text) or re.match(r"^\*\*[^*]{2,120}\*\*\s*$", text))


def _is_seed_block_heading(line: str) -> bool:
    text = str(line or "").strip()
    lower = text.lower()
    heading_like = _is_markdown_heading(text) or lower.startswith(("seed ", "seed-", "seed_", "seed:"))
    if not heading_like:
        return False
    return "seed" in lower and any(
        token in lower for token in ("list", "entities", "firms", "companies", "topics", "items", "files", "targets")
    )


def _seed_block_lines(instruction_text: str) -> list[str]:
    lines = str(instruction_text or "").splitlines()
    collecting = False
    collected: list[str] = []
    for raw_line in lines:
        stripped = raw_line.strip()
        if not stripped:
            continue
        if _is_seed_block_heading(stripped):
            collecting = True
            continue
        if not collecting:
            continue
        if stripped.startswith("---"):
            collecting = False
            continue
        if _is_markdown_heading(stripped):
            # Category headings inside a seed block are labels, not seed terms.
            continue
        lower = stripped.lower()
        if lower.startswith("(") and "expand beyond" in lower:
            continue
        if _contains_any(lower, ("aim for", "do not limit", "before scoring", "total before")):
            continue
        if "," in stripped or ";" in stripped or re.search(r"\band\b", stripped, re.I):
            collected.append(stripped)
    return collected


def _coverage_expectations(instruction_text: str) -> list[dict[str, object]]:
    expectations: list[dict[str, object]] = []
    for line in _lines(instruction_text):
        matches: list[tuple[int, int | None, str]] = []
        for pattern in (_COVERAGE_RANGE_RE, _COVERAGE_RANGE_TOTAL_RE):
            for match in pattern.finditer(line):
                low = int(match.group(1))
                high = int(match.group(2))
                if low > high:
                    low, high = high, low
                matches.append((low, high, match.group(3).lower()))
        for match in _COVERAGE_RANGE_BARE_TOTAL_RE.finditer(line):
            low = int(match.group(1))
            high = int(match.group(2))
            if low > high:
                low, high = high, low
            matches.append((low, high, "items"))
        for match in _COVERAGE_MIN_RE.finditer(line):
            matches.append((int(match.group(1)), None, match.group(2).lower()))
        if re.search(r"\b(?:aim|target|score|screen|cover|include|research|total)\b", line, re.I):
            for match in _COVERAGE_PLUS_RE.finditer(line):
                matches.append((int(match.group(1)), None, match.group(2).lower()))
        for minimum, maximum, entity in matches:
            expectations.append(
                {
                    "minimum": minimum,
                    "maximum": maximum,
                    "entity": entity,
                    "line": _redact_text(line),
                }
            )
    seen: set[tuple[int, int | None, str, str]] = set()
    unique: list[dict[str, object]] = []
    for item in expectations:
        key = (
            int(item.get("minimum") or 0),
            int(item["maximum"]) if item.get("maximum") is not None else None,
            str(item.get("entity") or ""),
            str(item.get("line") or ""),
        )
        if key in seen:
            continue
        seen.add(key)
        unique.append(item)
    return unique


def build_constraint_ledger(
    *,
    instruction: str,
    worker: dict[str, object],
    run_id: str,
    created_at: float | None = None,
) -> dict[str, object]:
    """Build a generic, conservative ledger from the user's instruction and structured inputs."""

    instruction_text = str(instruction or "")
    source_lines: list[str] = []
    date_lines: list[str] = []
    auth_lines: list[str] = []
    scope_lines: list[str] = []
    exclusion_lines: list[str] = []
    required_output_lines: list[str] = []
    forbidden_output_lines: list[str] = []
    seed_lines: list[str] = []
    memory_write_mode_off = _memory_write_mode_is_off(instruction_text)

    for line in _lines(instruction_text):
        lower = line.lower()
        if _line_has_date_constraint_context(line):
            date_lines.append(line)
        if _contains_any(lower, ("source", "citation", "primary", "official", "public source", "web source", "evidence")):
            source_lines.append(line)
        if _contains_any(lower, ("auth", "login", "signed", "credential", "token", "session", "account")):
            auth_lines.append(line)
        if _contains_any(lower, ("only", "scope", "in scope", "out of scope", "do not", "must", "exclude", "include", "flag", "local", "cloud")):
            scope_lines.append(line)
        if _contains_any(lower, ("exclude", "do not exclude", "flag", "forbid", "forbidden", "must not", "not allowed")):
            exclusion_lines.append(line)
        if _contains_any(lower, ("seed", "seed file", "seed entity", "input file", "uploaded file")) and not _is_seed_block_heading(line):
            seed_lines.append(line)
        if _line_has_required_output_context(line) and not (
            memory_write_mode_off and _line_is_memory_proposal_output_context(line)
        ):
            required_output_lines.append(line)
        if _line_forbids_output(line):
            forbidden_output_lines.extend(_forbidden_output_fragments(line))

    def unique_redacted(values: list[str]) -> list[str]:
        return list(dict.fromkeys(_redact_text(item) for item in values))

    seeds = unique_redacted([*seed_lines, *_seed_block_lines(instruction_text), *_bootstrap_seed_files(worker)])
    do_not_widen_or_soften = bool(
        date_lines
        or source_lines
        or any(
            _contains_any(line.lower(), ("only", "do not", "must", "no later than", "through", "until", "scope", "flag"))
            for line in _lines(instruction_text)
        )
    )
    return {
        "schema": "glasshive.run.constraint-ledger.v1",
        "run_id": run_id,
        "created_at": created_at if created_at is not None else _now_epoch(),
        "worker": {
            "worker_id": str(worker.get("worker_id") or ""),
            "profile": str(worker.get("profile") or ""),
            "execution_mode": str(worker.get("execution_mode") or ""),
            "runtime": str(worker.get("runtime") or ""),
            "backend": _derived_legacy_backend(worker),
        },
        "original_request": _redact_text(instruction_text),
        "constraints": {
            "date": unique_redacted(date_lines),
            "source": unique_redacted(source_lines),
            "auth": unique_redacted(auth_lines),
            "scope": unique_redacted(scope_lines),
            "exclusion_or_flag": unique_redacted(exclusion_lines),
        },
        "outputs": {
            "required": unique_redacted(required_output_lines),
            "forbidden": unique_redacted(forbidden_output_lines),
            "format_expectations": _required_output_formats(required_output_lines),
            "forbidden_format_expectations": _output_formats(forbidden_output_lines),
        },
        "seed_entities_or_files": seeds,
        "coverage_expectations": _coverage_expectations(instruction_text),
        "do_not_widen_or_soften": do_not_widen_or_soften,
    }


def _output_formats(lines: list[str]) -> list[str]:
    found: list[str] = []
    for suffix in _OUTPUT_FORMAT_SUFFIXES:
        if any(re.search(rf"\b{re.escape(suffix)}\b", line, re.I) for line in lines):
            found.append(suffix)
    return found


def _line_forbids_output(line: str) -> bool:
    return bool(_forbidden_output_fragments(line))


def _forbidden_output_fragments(line: str) -> list[str]:
    text = str(line or "")
    fragments: list[str] = []
    for match in _OUTPUT_FORBIDDEN_CLAUSE_RE.finditer(text):
        fragment = match.group(0).strip(" \t\n\r:;,.")
        if not fragment:
            continue
        if _OUTPUT_FORBIDDEN_RE.search(fragment) or _output_formats([fragment]):
            fragments.append(fragment)
    if not fragments:
        fallback = _OUTPUT_FORBIDDEN_RE.search(text)
        if fallback:
            fragments.append(fallback.group(0).strip(" \t\n\r:;,."))
    return list(dict.fromkeys(fragments))


def _line_has_required_output_context(line: str) -> bool:
    text = str(line or "")
    required_fragment = _OUTPUT_FORBIDDEN_CLAUSE_RE.sub(" ", text)
    return bool(_OUTPUT_ACTION_RE.search(required_fragment))


def _line_is_final_answer_only_context(line: str) -> bool:
    text = str(line or "")
    required_fragment = _OUTPUT_FORBIDDEN_CLAUSE_RE.sub(" ", text)
    if not _FINAL_ANSWER_ONLY_CONTEXT_RE.search(required_fragment):
        return False
    return not bool(_FILE_DELIVERABLE_CONTEXT_RE.search(required_fragment))


def _format_mentions(line: str, suffix: str) -> Iterable[re.Match[str]]:
    return re.finditer(rf"\b{re.escape(suffix)}\b", str(line or ""), re.I)


def _format_looks_like_filename_extension(line: str, match: re.Match[str]) -> bool:
    text = str(line or "")
    if match.start() == 0 or text[match.start() - 1] != ".":
        return False
    prefix = text[max(0, match.start() - 100) : match.start()]
    if not re.search(r"(?:^|[\s\"'`(<\[])[A-Za-z0-9][A-Za-z0-9._-]{0,80}\.$", prefix):
        return False
    stem_prefix = prefix[:-1]
    clause = re.split(r"[;\n]|(?<=[.!?])\s+", stem_prefix)[-1]
    direct_name = re.search(r"\b(named|called|as|to)\s+[A-Za-z0-9][A-Za-z0-9._-]{0,80}$", clause, re.I)
    output_action = re.search(
        r"\b(create|generate|produce|write|save|export|deliver|return|output|build|prepare|render)\b",
        clause,
        re.I,
    )
    input_descriptor = re.search(r"\b(attached|uploaded|input|source|provided|existing|original|reference)\b", clause, re.I)
    input_named_file = re.search(
        r"\b(attached|uploaded|input|source|provided|existing|original|reference)\b"
        r"[^;\n]{0,80}\b(file|attachment|document|spreadsheet|workbook)\b"
        r"[^;\n]{0,35}\b(?:is\s+)?(?:named|called)\s+[A-Za-z0-9][A-Za-z0-9._-]{0,80}$",
        clause,
        re.I,
    )
    if input_named_file:
        return False
    return bool(direct_name or (output_action and not input_descriptor))


def _format_looks_like_input_filename_extension(line: str, match: re.Match[str]) -> bool:
    text = str(line or "")
    if match.start() == 0 or text[match.start() - 1] != ".":
        return False
    prefix = text[max(0, match.start() - 120) : match.start()]
    if not re.search(r"(?:^|[\s\"'`(<\[])[A-Za-z0-9][A-Za-z0-9._-]{0,100}\.$", prefix):
        return False
    stem_prefix = prefix[:-1]
    clause = re.split(r"[;\n]|(?<=[.!?])\s+", stem_prefix)[-1]
    return bool(
        re.search(
            r"\b(attached|uploaded|input|source|provided|existing|original|reference)\b"
            r"[^;\n]{0,80}\b(file|attachment|document|spreadsheet|workbook)\b"
            r"[^;\n]{0,35}\b(?:is\s+)?(?:named|called)\s+[A-Za-z0-9][A-Za-z0-9._-]{0,100}$",
            clause,
            re.I,
        )
    )


def _is_input_format_mention(line: str, match: re.Match[str]) -> bool:
    line_text = str(line or "")
    before = line_text[max(0, match.start() - 60) : match.start()]
    after = line_text[match.end() : min(len(line_text), match.end() + 60)]
    nearby = f"{before} {after}"
    immediate = f"{before[-45:]} {after[:20]}"
    suffix = str(match.group(0) or "").lower()
    if match.start() > 0 and line_text[match.start() - 1] == ".":
        path_prefix = line_text[max(0, match.start() - 120) : match.start()]
        if re.search(
            r"(?:^|[\s\"'`(<\[])(?:uploads?|inputs?|sources?|provided|existing|original|reference|workspace)"
            r"[/\\][^\s\"'`),;]*\.$",
            path_prefix,
            re.I,
        ):
            return True
    if _format_looks_like_input_filename_extension(line_text, match):
        return True
    if _format_looks_like_filename_extension(line_text, match):
        return False
    if match.start() > 0 and line_text[match.start() - 1] == "/" and re.search(
        r"\b(application|image|audio|video|text|multipart)\s*/\s*$",
        before,
        re.I,
    ):
        return True
    if suffix:
        if suffix == "pdf" and (before.rstrip().endswith("%") or re.search(r"\bsignature\b", nearby, re.I)):
            return True
    if re.search(r"^\s*(bytes?|content|data)\b", after, re.I) and not re.search(
        r"\b(create|generate|produce|write|save|export|deliver|return|build|prepare|compose|draft|render)\b[^.\n;]{0,35}$",
        before,
        re.I,
    ):
        return True
    if re.search(
        r"\b(?:of|from|using|with|based\s+on)?\s*(?:the\s+)?"
        r"(attached|uploaded|input|source|provided|existing|original|reference|embedded)\s*$",
        before[-70:],
        re.I,
    ):
        return True
    if suffix == "pdf" and re.search(r"\b(bytes?|signature|header|first[-\s]*byte)\b", nearby, re.I) and re.search(
        r"\b(attached|uploaded|input|source|provided|existing|original|workspace)\b", nearby, re.I
    ):
        return True
    if re.search(r"\b(read|inspect|analy[sz]e|parse|verify|use)\b[^.\n;]{0,45}$", before, re.I) and re.search(
        r"^\s*(?:from|in|inside|under|within|at)\b|\buploads?/|\bworkspace\b",
        after,
        re.I,
    ):
        return True
    if re.search(
        r"\b(from|using|with|based\s+on)\b[^.\n;]{0,45}\b(attached|uploaded|input|source|provided|existing|original)?\s*$",
        before,
        re.I,
    ):
        return True
    if re.search(
        r"\b(create|generate|produce|write|save|export|deliver|return|build|prepare|compose|draft|render)\b[^.\n;]{0,35}$",
        before,
        re.I,
    ):
        return False
    if re.search(r"^\s*(?:artifact|artifacts|file|files|report|reports|deliverable|deliverables)\b", after, re.I):
        return False
    if re.search(r"\b(attached|uploaded|input|source|provided|existing|original|reference|embedded)\b", immediate, re.I):
        return True
    if re.search(r"\b(output|outputs|artifact|artifacts|deliverable|deliverables|report|reports)\b", nearby, re.I):
        return False
    return bool(_INPUT_FORMAT_CONTEXT_RE.search(nearby))


def _is_forbidden_format_mention(line: str, match: re.Match[str]) -> bool:
    start = max(0, match.start() - 80)
    end = min(len(str(line or "")), match.end() + 80)
    return bool(_OUTPUT_FORBIDDEN_RE.search(str(line or "")[start:end]))


def _required_output_formats(lines: list[str]) -> list[str]:
    found: list[str] = []
    for line in lines:
        if not _line_has_required_output_context(line):
            continue
        scan_line = _OUTPUT_FORBIDDEN_CLAUSE_RE.sub(" ", str(line or ""))
        for suffix in _OUTPUT_FORMAT_SUFFIXES:
            for match in _format_mentions(scan_line, suffix):
                if _is_forbidden_format_mention(scan_line, match) or _is_input_format_mention(scan_line, match):
                    continue
                found.append(suffix)
                break
    return list(dict.fromkeys(found))


def write_constraint_ledger(workspace_dir: Path | str, ledger: dict[str, object], run_id: str) -> Path:
    run_dir = Path(workspace_dir) / RUN_EVIDENCE_DIR_NAME
    latest_path = run_dir / "constraint-ledger.json"
    per_run_path = run_dir / "runs" / run_id / "constraint-ledger.json"
    latest_path.parent.mkdir(parents=True, exist_ok=True)
    per_run_path.parent.mkdir(parents=True, exist_ok=True)
    body = json.dumps(ledger, indent=2, sort_keys=True)
    latest_path.write_text(body + "\n")
    per_run_path.write_text(body + "\n")
    return latest_path


def _walk_values(value: object, prefix: str = "") -> Iterable[tuple[str, str]]:
    if isinstance(value, dict):
        for key, nested in value.items():
            next_prefix = f"{prefix}.{key}" if prefix else str(key)
            yield from _walk_values(nested, next_prefix)
    elif isinstance(value, list):
        for index, nested in enumerate(value):
            yield from _walk_values(nested, f"{prefix}[{index}]")
    else:
        yield prefix or "$", str(value or "")


def check_content_hygiene(value: object) -> dict[str, object]:
    issues: list[dict[str, str]] = []
    for path, raw_text in _walk_values(value):
        text = " ".join(raw_text.split())
        if not text:
            continue
        for pattern, reason in _CONTENT_HYGIENE_PATTERNS:
            match = pattern.search(text)
            if not match:
                continue
            if _CONTENT_HYGIENE_NEGATIVE_CONTEXT_RE.search(text[max(0, match.start() - 120) : match.start()]):
                continue
            issues.append(
                {
                    "path": path,
                    "reason": reason,
                    "match": match.group(0),
                    "text": _redact_text(text[:240]),
                }
            )
            break
        structured_path = bool(re.search(r"(\.csv|\.tsv|\.jsonl?|[\].][A-Za-z0-9_-]+)$", path, re.I)) and not re.search(
            r"\.(md|txt)$", path, re.I
        )
        if (
            len(text) > 900
            and structured_path
            and not re.search(
                r"(source|evidence|notes|summary|description|rationale|transcript|narrative|profile|case|pain|"
                r"synthesis|assessment|sequence|angle|hook|recommendation|positioning|outreach|finding)",
                path,
                re.I,
            )
        ):
            issues.append(
                {
                    "path": path,
                    "reason": "overlong structured field",
                    "match": "overlong structured field",
                    "text": _redact_text(text[:240]),
                }
            )
    return {
        "status": "warn" if issues else "pass",
        "failure_count": len(issues),
        "issues": issues[:100],
    }


def _relative_path(path: Path, root: Path) -> str:
    try:
        return path.relative_to(root).as_posix()
    except ValueError:
        return path.name


def _document_validation(path: Path) -> dict[str, object]:
    suffix = path.suffix.lower()
    if suffix not in PROFESSIONAL_ARTIFACT_EXTENSIONS:
        return {"applicable": False, "valid": None, "method": "not_applicable"}
    validation: dict[str, object] = {
        "applicable": True,
        "valid": is_valid_professional_artifact(path),
        "method": "structural",
    }
    if suffix == ".pdf":
        validation["method"] = "pdf_header"
        validation["pdfinfo"] = _pdfinfo_validation(path)
        validation["render_sample"] = _pdf_render_validation(path)
        return validation
    if suffix in OOXML_ARTIFACT_MARKERS:
        validation["method"] = "ooxml_marker"
        validation["required_marker"] = OOXML_ARTIFACT_MARKERS[suffix]
        validation["library_check"] = _ooxml_library_validation(path, suffix)
        return validation
    if suffix in ODF_ARTIFACT_EXTENSIONS:
        validation["method"] = "odf_marker"
        return validation
    if suffix in OLE_ARTIFACT_EXTENSIONS:
        validation["method"] = "ole_magic"
        return validation
    return validation


def _pdfinfo_validation(path: Path) -> dict[str, object]:
    binary = shutil.which("pdfinfo")
    if not binary:
        return {"status": "unavailable"}
    try:
        completed = subprocess.run(
            [binary, str(path)],
            check=False,
            capture_output=True,
            text=True,
            timeout=10,
        )
    except Exception as exc:  # pragma: no cover - defensive evidence path
        return {"status": "error", "error": _redact_text(exc)}
    output = completed.stdout or completed.stderr or ""
    pages_match = re.search(r"^Pages:\s*(\d+)", output, re.MULTILINE)
    return {
        "status": "pass" if completed.returncode == 0 else "warn",
        "returncode": completed.returncode,
        "pages": int(pages_match.group(1)) if pages_match else None,
        "tail": _redact_text(output, max_chars=1000),
    }


def _pdf_render_validation(path: Path) -> dict[str, object]:
    binary = shutil.which("pdftoppm")
    if not binary:
        return {"status": "unavailable"}
    try:
        with tempfile.TemporaryDirectory(prefix="glasshive-pdf-render-") as temp_dir:
            output_prefix = Path(temp_dir) / "sample"
            sample_path = output_prefix.with_suffix(".png")
            completed = subprocess.run(
                [binary, "-singlefile", "-png", "-f", "1", "-l", "1", str(path), str(output_prefix)],
                check=False,
                capture_output=True,
                text=True,
                timeout=20,
            )
            exists = sample_path.exists()
    except Exception as exc:  # pragma: no cover - defensive evidence path
        return {"status": "error", "error": _redact_text(exc)}
    return {
        "status": "pass" if completed.returncode == 0 and exists else "warn",
        "returncode": completed.returncode,
        "sample_created": exists,
        "tail": _redact_text((completed.stdout or completed.stderr or ""), max_chars=1000),
    }


def _ooxml_library_validation(path: Path, suffix: str) -> dict[str, object]:
    if suffix == ".xlsx":
        try:
            import openpyxl  # type: ignore
        except Exception:
            return {"status": "unavailable", "library": "openpyxl"}
        try:
            workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
            sheet_count = len(workbook.sheetnames)
            workbook.close()
            return {"status": "pass", "library": "openpyxl", "sheet_count": sheet_count}
        except Exception as exc:
            return {"status": "warn", "library": "openpyxl", "error": _redact_text(exc)}
    if suffix == ".docx":
        try:
            import docx  # type: ignore
        except Exception:
            return {"status": "unavailable", "library": "python-docx"}
        try:
            document = docx.Document(str(path))
            return {"status": "pass", "library": "python-docx", "paragraph_count": len(document.paragraphs)}
        except Exception as exc:
            return {"status": "warn", "library": "python-docx", "error": _redact_text(exc)}
    if suffix == ".pptx":
        try:
            import pptx  # type: ignore
        except Exception:
            return {"status": "unavailable", "library": "python-pptx"}
        try:
            presentation = pptx.Presentation(str(path))
            return {"status": "pass", "library": "python-pptx", "slide_count": len(presentation.slides)}
        except Exception as exc:
            return {"status": "warn", "library": "python-pptx", "error": _redact_text(exc)}
    return {"status": "not_applicable"}


def _scheduled_prompt_private_scratchpad_dir(workspace_dir: Path) -> Path | None:
    snapshot_path = workspace_dir / "scheduled-prompt" / "variable-snapshot.json"
    if not snapshot_path.is_file():
        return None
    try:
        payload = json.loads(snapshot_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return None
    items = payload.get("items") if isinstance(payload, dict) else None
    if not isinstance(items, list):
        return None
    for item in items:
        if not isinstance(item, dict):
            continue
        marker = str(item.get("placeholder") or item.get("wrapper") or "").strip()
        if marker not in {
            "local.viventium.my_folder",
            "local.viventium.local_machine_glasshive.my_folder",
        }:
            continue
        value = str(item.get("value") or "").strip()
        if not value:
            continue
        path = Path(value).expanduser()
        try:
            if path.is_dir():
                return path
        except OSError:
            return None
    return None


def _private_scratchpad_artifact_paths(
    workspace_dir: Path,
    *,
    started_at: float | None = None,
    ended_at: float | None = None,
) -> list[Path]:
    root = _scheduled_prompt_private_scratchpad_dir(workspace_dir)
    if root is None:
        return []
    lower_bound = (started_at - 300) if started_at is not None else None
    upper_bound = (ended_at + 300) if ended_at is not None else None
    candidates: list[Path] = []
    for path in sorted(root.iterdir()):
        if not path.is_file():
            continue
        if not (
            re.fullmatch(r"\d{12}\.md", path.name)
            or re.fullmatch(r"memory-proposals-\d{12}\.json", path.name)
        ):
            continue
        try:
            mtime = path.stat().st_mtime
        except OSError:
            continue
        if lower_bound is not None and mtime < lower_bound:
            continue
        if upper_bound is not None and mtime > upper_bound:
            continue
        candidates.append(path)
    return candidates[-20:]


def _artifact_inventory(
    workspace_dir: Path,
    *,
    started_at: float | None = None,
    ended_at: float | None = None,
) -> dict[str, object]:
    items: list[dict[str, object]] = []
    for path in candidate_artifact_paths({"workspace_dir": str(workspace_dir)}, max_entries=200):
        try:
            stat = path.stat()
        except OSError:
            continue
        items.append(
            {
                "path": _relative_path(path, workspace_dir),
                "bytes": stat.st_size,
                "mtime": stat.st_mtime,
                "suffix": path.suffix.lower(),
                "document_validation": _document_validation(path),
            }
        )
        if path.suffix.lower() in _HTML_FILE_SUFFIXES:
            items[-1]["html_validation"] = _html_browser_validation(path)
    for path in _private_scratchpad_artifact_paths(workspace_dir, started_at=started_at, ended_at=ended_at):
        try:
            stat = path.stat()
        except OSError:
            continue
        items.append(
            {
                "path": f"private-scratchpad/{path.name}",
                "bytes": stat.st_size,
                "mtime": stat.st_mtime,
                "suffix": path.suffix.lower(),
                "source": "scheduled_prompt_private_scratchpad",
                "document_validation": _document_validation(path),
            }
        )
    return {
        "count": len(items),
        "items": items,
    }


def _html_browser_validation(path: Path) -> dict[str, object]:
    enabled = os.environ.get("GLASSHIVE_EVIDENCE_HTML_PLAYWRIGHT", "").strip().lower()
    if enabled not in {"1", "true", "yes", "on"}:
        return {
            "status": "unavailable",
            "reason": "disabled; set GLASSHIVE_EVIDENCE_HTML_PLAYWRIGHT=1",
        }
    node = shutil.which("node")
    if not node:
        return {"status": "unavailable", "reason": "node not found"}
    script = r"""
const target = process.argv[1];
(async () => {
  let chromium;
  try {
    chromium = require("playwright").chromium;
  } catch (error) {
    console.log(JSON.stringify({ status: "unavailable", reason: "playwright not installed" }));
    return;
  }
  const browser = await chromium.launch({ headless: true });
  const page = await browser.newPage();
  const consoleMessages = [];
  page.on("console", (message) => {
    if (["error", "warning"].includes(message.type())) {
      consoleMessages.push({ type: message.type(), text: message.text().slice(0, 240) });
    }
  });
  page.on("pageerror", (error) => {
    consoleMessages.push({ type: "pageerror", text: String(error).slice(0, 240) });
  });
  await page.goto(target, { waitUntil: "load", timeout: 10000 });
  const bodyText = await page.locator("body").innerText({ timeout: 5000 }).catch(() => "");
  const title = await page.title();
  await browser.close();
  console.log(JSON.stringify({
    status: consoleMessages.some((item) => item.type === "error" || item.type === "pageerror") ? "warn" : "pass",
    title,
    body_chars: bodyText.length,
    console: consoleMessages.slice(0, 20),
  }));
})().catch((error) => {
  console.log(JSON.stringify({ status: "warn", error: String(error).slice(0, 500) }));
  process.exitCode = 1;
});
"""
    try:
        completed = subprocess.run(
            [node, "-e", script, path.resolve().as_uri()],
            check=False,
            capture_output=True,
            text=True,
            timeout=20,
        )
    except Exception as exc:  # pragma: no cover - defensive evidence path
        return {"status": "error", "error": _redact_text(exc)}
    raw = (completed.stdout or completed.stderr or "").strip()
    try:
        payload = json.loads(raw.splitlines()[-1]) if raw else {}
    except json.JSONDecodeError:
        payload = {"status": "warn", "tail": _redact_text(raw, max_chars=1000)}
    payload.setdefault("returncode", completed.returncode)
    return payload


def _visual_render_summary(artifacts: dict[str, object]) -> dict[str, object]:
    items = artifacts.get("items") if isinstance(artifacts, dict) else []
    evidence: list[dict[str, object]] = []
    if not isinstance(items, list):
        return {"items": evidence, "status": "not_applicable"}
    for item in items:
        if not isinstance(item, dict):
            continue
        path = str(item.get("path") or "")
        document_validation = item.get("document_validation")
        if isinstance(document_validation, dict):
            render_sample = document_validation.get("render_sample")
            if isinstance(render_sample, dict) and render_sample.get("status") != "unavailable":
                evidence.append({"path": path, "kind": "pdf_render_sample", **render_sample})
        html_validation = item.get("html_validation")
        if isinstance(html_validation, dict):
            evidence.append({"path": path, "kind": "html_browser_smoke", **html_validation})
    if not evidence:
        return {"items": evidence, "status": "not_available"}
    if any(str(item.get("status") or "").lower() in {"warn", "error"} for item in evidence):
        return {"items": evidence, "status": "warn"}
    if any(str(item.get("status") or "").lower() == "pass" for item in evidence):
        return {"items": evidence, "status": "pass"}
    return {"items": evidence, "status": "unavailable"}


def _text_artifact_payload(workspace_dir: Path) -> dict[str, str]:
    payload: dict[str, str] = {}
    for path in candidate_artifact_paths({"workspace_dir": str(workspace_dir)}, max_entries=100):
        if path.suffix.lower() not in _CONTENT_HYGIENE_TEXT_SUFFIXES:
            continue
        try:
            if path.stat().st_size > _MAX_TEXT_SCAN_BYTES:
                continue
            rel = _relative_path(path, workspace_dir)
            if _is_raw_support_capture_path(rel):
                continue
            text = path.read_text(encoding="utf-8", errors="ignore")
            suffix = path.suffix.lower()
            if suffix in {".csv", ".tsv"}:
                delimiter = "\t" if suffix == ".tsv" else ","
                try:
                    reader = csv.reader(text.splitlines(), delimiter=delimiter)
                    for row_no, row in enumerate(reader, start=1):
                        if row_no > 1000:
                            break
                        for col_no, cell in enumerate(row, start=1):
                            if col_no > 100:
                                break
                            if str(cell or "").strip():
                                payload[f"{rel}:R{row_no}C{col_no}"] = cell
                    continue
                except csv.Error:
                    payload[rel] = text
                    continue
            if suffix == ".json":
                try:
                    parsed = json.loads(text)
                except json.JSONDecodeError:
                    payload[rel] = text
                    continue
                for index, (nested_path, value) in enumerate(_walk_values(parsed)):
                    if index >= 1000:
                        break
                    payload[f"{rel}:{nested_path}"] = value
                continue
            if suffix == ".jsonl":
                parsed_any = False
                for line_no, line in enumerate(text.splitlines(), start=1):
                    if line_no > 1000:
                        break
                    if not line.strip():
                        continue
                    try:
                        parsed = json.loads(line)
                    except json.JSONDecodeError:
                        payload[f"{rel}:{line_no}"] = line
                        continue
                    parsed_any = True
                    for nested_path, value in _walk_values(parsed):
                        payload[f"{rel}:{line_no}:{nested_path}"] = value
                if parsed_any:
                    continue
            payload[rel] = text
        except OSError:
            continue
    return payload


def _non_empty_csv_rows(path: Path, *, delimiter: str) -> int:
    try:
        text = path.read_text(encoding="utf-8", errors="ignore")
        reader = csv.reader(text.splitlines(), delimiter=delimiter)
        rows = [row for row in reader if any(str(cell or "").strip() for cell in row)]
    except (OSError, csv.Error):
        return 0
    return max(0, len(rows) - 1) if rows else 0


def _json_collection_count(path: Path) -> int:
    try:
        parsed = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
    except (OSError, json.JSONDecodeError):
        return 0
    if isinstance(parsed, list):
        return len(parsed)
    if isinstance(parsed, dict):
        counts = [len(value) for value in parsed.values() if isinstance(value, list)]
        return max(counts) if counts else 0
    return 0


def _xlsx_table_counts(path: Path) -> list[dict[str, object]]:
    try:
        import openpyxl  # type: ignore
    except Exception:
        return []
    try:
        workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    except Exception:
        return []
    counts: list[dict[str, object]] = []
    try:
        for sheet in workbook.worksheets:
            non_empty_rows = 0
            for row in sheet.iter_rows(values_only=True):
                if any(str(cell or "").strip() for cell in row):
                    non_empty_rows += 1
            counts.append({"sheet": sheet.title, "count": max(0, non_empty_rows - 1)})
    finally:
        workbook.close()
    return counts


def _markdown_table_count(path: Path) -> int:
    try:
        lines = path.read_text(encoding="utf-8", errors="ignore").splitlines()
    except OSError:
        return 0
    return _markdown_table_count_from_lines(lines)


def _markdown_table_count_from_lines(lines: Iterable[str]) -> int:
    table_rows = [
        line
        for line in lines
        if "|" in line and not re.fullmatch(r"\s*\|?\s*:?-{3,}:?\s*(\|\s*:?-{3,}:?\s*)+\|?\s*", line)
    ]
    if len(table_rows) < 2:
        return 0
    return max(0, len(table_rows) - 1)


def _html_table_count(path: Path) -> int:
    try:
        text = path.read_text(encoding="utf-8", errors="ignore")
    except OSError:
        return 0
    rows = len(re.findall(r"<tr\b", text, re.I))
    return max(0, rows - 1)


def _final_output_coverage_counts(output_text: str) -> list[dict[str, object]]:
    text = str(output_text or "")
    if not text.strip():
        return []
    lines = text.splitlines()
    counts: list[dict[str, object]] = []
    list_count = sum(
        1
        for line in lines
        if re.match(r"^\s*(?:[-*+]\s+|\d+[.)]\s+)", line)
    )
    if list_count:
        counts.append({"path": "final_output", "kind": "final_output_list_items", "count": list_count})
    table_count = _markdown_table_count_from_lines(lines)
    if table_count:
        counts.append({"path": "final_output", "kind": "final_output_table_rows", "count": table_count})
    return counts


def _artifact_coverage_counts(workspace_dir: Path, artifacts: dict[str, object], *, output_text: str = "") -> list[dict[str, object]]:
    items = artifacts.get("items") if isinstance(artifacts, dict) else []
    if not isinstance(items, list):
        return []
    counts: list[dict[str, object]] = []
    for item in items:
        if not isinstance(item, dict):
            continue
        rel_path = str(item.get("path") or "")
        if not rel_path or _is_support_artifact_path(rel_path):
            continue
        path = workspace_dir / rel_path
        suffix = path.suffix.lower()
        if suffix == ".csv":
            count = _non_empty_csv_rows(path, delimiter=",")
            if count:
                counts.append({"path": rel_path, "kind": "csv_rows", "count": count})
        elif suffix == ".tsv":
            count = _non_empty_csv_rows(path, delimiter="\t")
            if count:
                counts.append({"path": rel_path, "kind": "tsv_rows", "count": count})
        elif suffix == ".json":
            count = _json_collection_count(path)
            if count:
                counts.append({"path": rel_path, "kind": "json_items", "count": count})
        elif suffix == ".xlsx":
            for sheet_count in _xlsx_table_counts(path):
                counts.append({"path": rel_path, "kind": "xlsx_rows", **sheet_count})
        elif suffix == ".md":
            count = _markdown_table_count(path)
            if count:
                counts.append({"path": rel_path, "kind": "markdown_table_rows", "count": count})
        elif suffix in _HTML_FILE_SUFFIXES:
            count = _html_table_count(path)
            if count:
                counts.append({"path": rel_path, "kind": "html_table_rows", "count": count})
    counts.extend(_final_output_coverage_counts(output_text))
    return counts


def _coverage_compliance(
    workspace_dir: Path,
    artifacts: dict[str, object],
    ledger: dict[str, object] | None,
    *,
    output_text: str = "",
) -> dict[str, object]:
    expectations: list[dict[str, object]] = []
    if isinstance(ledger, dict) and isinstance(ledger.get("coverage_expectations"), list):
        expectations = [item for item in ledger["coverage_expectations"] if isinstance(item, dict)]
    counts = _artifact_coverage_counts(workspace_dir, artifacts, output_text=output_text)
    count_values = [int(item.get("count") or 0) for item in counts]
    max_count = max(count_values, default=0)
    matched_count = 0
    issues: list[dict[str, object]] = []
    warnings: list[dict[str, object]] = []
    for expectation in expectations:
        minimum = int(expectation.get("minimum") or 0)
        raw_maximum = expectation.get("maximum")
        maximum = int(raw_maximum) if raw_maximum is not None else None
        matching_counts = [
            value
            for value in count_values
            if (not minimum or value >= minimum) and (maximum is None or value <= maximum)
        ]
        if matching_counts:
            matched_count = max(matched_count, max(matching_counts))
            continue
        if minimum and not any(value >= minimum for value in count_values):
            payload = {
                "reason": "coverage count below requested minimum" if count_values else "coverage count could not be verified",
                "requested_minimum": minimum,
                "requested_maximum": maximum,
                "observed_max_count": max_count,
                "entity": str(expectation.get("entity") or ""),
                "line": _redact_text(expectation.get("line") or ""),
            }
            if count_values:
                issues.append(payload)
            else:
                warnings.append(payload)
        elif maximum is not None and any(value > maximum for value in count_values):
            warnings.append(
                {
                    "reason": "coverage count above requested maximum",
                    "requested_minimum": minimum,
                    "requested_maximum": maximum,
                    "observed_max_count": max_count,
                    "entity": str(expectation.get("entity") or ""),
                    "line": _redact_text(expectation.get("line") or ""),
                }
            )
    return {
        "status": "fail" if issues else "warn" if warnings else "pass" if expectations else "not_applicable",
        "expectations": expectations,
        "counts": counts[:50],
        "observed_max_count": max_count,
        "matched_count": matched_count or None,
        "issues": [*issues, *warnings][:100],
    }


def _iter_scan_files(workspace_dir: Path) -> Iterable[Path]:
    count = 0
    for path in workspace_dir.rglob("*"):
        if count >= _MAX_TEXT_SCAN_FILES:
            return
        if not path.is_file() or path.suffix.lower() not in _TEXT_FILE_SUFFIXES:
            continue
        try:
            rel = path.relative_to(workspace_dir)
        except ValueError:
            continue
        if _is_raw_support_capture_path(rel.as_posix()):
            continue
        lowered_parts = [part.lower() for part in rel.parts]
        if any(part in _SKIP_SCAN_DIRS for part in lowered_parts):
            continue
        if not is_user_deliverable_relative_path(rel) and rel.parts[0].lower() not in {"research", "planning", "specs", "notes"}:
            continue
        try:
            if path.stat().st_size > _MAX_TEXT_SCAN_BYTES:
                continue
        except OSError:
            continue
        count += 1
        yield path


def _extract_xlsx_text(path: Path) -> tuple[str, str | None]:
    try:
        import openpyxl  # type: ignore
    except Exception:
        return "", "openpyxl unavailable for xlsx constraint scan"
    values: list[str] = []
    try:
        workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    except Exception:
        return "", "xlsx text extraction failed"
    try:
        for sheet_name in workbook.sheetnames[:20]:
            worksheet = workbook[sheet_name]
            values.append(str(sheet_name))
            for row_index, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
                if row_index > 1000:
                    break
                for value in row[:100]:
                    text = str(value or "").strip()
                    if text:
                        values.append(text)
                    if len(values) >= 5000:
                        return "\n".join(values), None
    finally:
        workbook.close()
    return "\n".join(values), None


def _extract_docx_text(path: Path) -> tuple[str, str | None]:
    try:
        import docx  # type: ignore
    except Exception:
        return "", "python-docx unavailable for docx constraint scan"
    try:
        document = docx.Document(str(path))
    except Exception:
        return "", "docx text extraction failed"
    fragments: list[str] = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    for table in document.tables[:50]:
        for row in table.rows[:1000]:
            for cell in row.cells[:100]:
                if cell.text:
                    fragments.append(cell.text)
                if len(fragments) >= 5000:
                    return "\n".join(fragments), None
    return "\n".join(fragments), None


def _extract_pptx_text(path: Path) -> tuple[str, str | None]:
    try:
        import pptx  # type: ignore
    except Exception:
        return "", "python-pptx unavailable for pptx constraint scan"
    try:
        presentation = pptx.Presentation(str(path))
    except Exception:
        return "", "pptx text extraction failed"
    fragments: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        if slide_index > 100:
            break
        for shape in slide.shapes:
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                fragments.append(text)
            if len(fragments) >= 5000:
                return "\n".join(fragments), None
    return "\n".join(fragments), None


def _extract_pdf_text(path: Path) -> tuple[str, str | None]:
    binary = shutil.which("pdftotext")
    if not binary:
        return "", "pdftotext unavailable for pdf constraint scan"
    try:
        completed = subprocess.run(
            [binary, "-layout", str(path), "-"],
            check=False,
            capture_output=True,
            text=True,
            timeout=20,
        )
    except Exception:
        return "", "pdf text extraction failed"
    if completed.returncode != 0:
        return "", "pdf text extraction failed"
    return completed.stdout[:_MAX_TEXT_SCAN_BYTES], None


def _extract_artifact_text(path: Path) -> tuple[str, str | None]:
    suffix = path.suffix.lower()
    try:
        if path.stat().st_size > _MAX_TEXT_SCAN_BYTES * 4:
            return "", "binary artifact too large for constraint scan"
    except OSError:
        return "", "binary artifact stat failed for constraint scan"
    if suffix == ".xlsx":
        return _extract_xlsx_text(path)
    if suffix == ".docx":
        return _extract_docx_text(path)
    if suffix == ".pptx":
        return _extract_pptx_text(path)
    if suffix == ".pdf":
        return _extract_pdf_text(path)
    return "", None


def _constraint_scan_texts(workspace_dir: Path, output_text: str) -> tuple[list[tuple[str, str]], list[dict[str, str]]]:
    entries: list[tuple[str, str]] = []
    warnings: list[dict[str, str]] = []
    if str(output_text or "").strip():
        entries.append(("<final-output>", str(output_text or "")))
    seen_paths: set[Path] = set()
    for path in _iter_scan_files(workspace_dir):
        seen_paths.add(path.resolve())
        try:
            entries.append((_relative_path(path, workspace_dir), path.read_text(encoding="utf-8", errors="ignore")))
        except OSError:
            continue
    for path in candidate_artifact_paths({"workspace_dir": str(workspace_dir)}, max_entries=100):
        if path.resolve() in seen_paths:
            continue
        if path.suffix.lower() not in _CONSTRAINT_BINARY_TEXT_SUFFIXES:
            continue
        text, extraction_warning = _extract_artifact_text(path)
        rel = _relative_path(path, workspace_dir)
        if extraction_warning:
            warnings.append(
                {
                    "path": rel,
                    "reason": "constraint text extraction unavailable",
                    "text": _redact_text(extraction_warning, max_chars=160),
                }
            )
        if text.strip():
            entries.append((rel, text[:_MAX_TEXT_SCAN_BYTES]))
    return entries, warnings


def _extract_latest_date_limit(ledger: dict[str, object]) -> tuple[int, int] | None:
    fragments: list[str] = []
    request = str(ledger.get("original_request") or "")
    fragments.append(request)
    constraints = ledger.get("constraints")
    if isinstance(constraints, dict):
        for value in constraints.values():
            if isinstance(value, list):
                fragments.extend(str(item) for item in value)
    limit: tuple[int, int] | None = None
    for fragment in fragments:
        for match in _DATE_LIMIT_RE.finditer(fragment):
            month = _MONTHS.get(match.group(1).lower()[:3]) or _MONTHS.get(match.group(1).lower())
            year = int(match.group(2))
            if month:
                candidate = (year, month)
                if limit is None or candidate > limit:
                    limit = candidate
    return limit


def _month_year_mentions(text: str) -> Iterable[tuple[int, int, str]]:
    for match in _MONTH_YEAR_RE.finditer(text):
        month = _MONTHS.get(match.group(1).lower()[:3]) or _MONTHS.get(match.group(1).lower())
        if month:
            yield int(match.group(2)), month, match.group(0)
    for match in _ISO_YEAR_MONTH_RE.finditer(text):
        yield int(match.group(1)), int(match.group(2)), match.group(0)


_REJECTED_OR_OUT_OF_SCOPE_RE = re.compile(
    r"\b("
    r"rejected|discarded|out[- ]of[- ]scope|outside scope|out[- ]of[- ]window|outside\b.{0,80}\bwindow|"
    r"excluded(?: from (?:evidence|facts|scoring|deliverables|factual record))?|not used|do not use|"
    r"not usable|not counted|not relied upon|not supporting|does not support|"
    r"removed from (?:facts|scoring|deliverables)"
    r")\b",
    re.I,
)
_ACCESS_TIMESTAMP_ONLY_RE = re.compile(
    r"\b(?:access(?:ed)?|retriev(?:ed|al))(?:\s*/\s*(?:access(?:ed)?|retriev(?:ed|al)))?\s+(?:date|timestamp|at|on|=|:|\d{4}-\d{2}|[A-Z][a-z]+ \d{4})\b"
    r"|\b(?:date|timestamp)\s+(?:of\s+)?(?:access|retrieval)\b",
    re.I,
)
_SOURCE_PUBLICATION_DATE_RE = re.compile(r"\b(?:published|publication|dated|released|filed|announced)\b", re.I)
_CONSTRAINT_PLANNING_DIRS = {"planning", "specs"}
_CONSTRAINT_SUPPORT_DIRS = _SUPPORT_ARTIFACT_DIRS
_CONSTRAINT_INTERNAL_NAMES = {"prompt", "subagent", "delegation", "constraint", "ledger"}
_CONSTRAINT_SUPPORT_NAMES = {"spec", "plan", "prompt", "subagent", "delegation", "constraint", "ledger"}
_SOFTENING_RE = re.compile(r"\b(wherever possible|if possible|best effort|when available|roughly|approximately)\b", re.I)
_CONSTRAINT_CONTEXT_RE = re.compile(
    r"\b(constraint|must|only|through|until|no later than|required|forbidden|exclude|flag)\b",
    re.I,
)


def _line_for_span(text: str, start: int, end: int) -> str:
    line_start = text.rfind("\n", 0, start) + 1
    line_end = text.find("\n", end)
    if line_end == -1:
        line_end = len(text)
    return text[line_start:line_end]


def _month_year_mentions_with_context(text: str) -> Iterable[tuple[int, int, str, str]]:
    for match in _MONTH_YEAR_RE.finditer(text):
        month = _MONTHS.get(match.group(1).lower()[:3]) or _MONTHS.get(match.group(1).lower())
        if month:
            yield int(match.group(2)), month, match.group(0), _line_for_span(text, match.start(), match.end())
    for match in _ISO_YEAR_MONTH_RE.finditer(text):
        yield int(match.group(1)), int(match.group(2)), match.group(0), _line_for_span(text, match.start(), match.end())


def _is_rejected_or_out_of_scope_line(line: str) -> bool:
    return bool(_REJECTED_OR_OUT_OF_SCOPE_RE.search(line))


def _is_access_timestamp_only_line(line: str) -> bool:
    text = str(line or "")
    return bool(_ACCESS_TIMESTAMP_ONLY_RE.search(text)) and not bool(_SOURCE_PUBLICATION_DATE_RE.search(text))


def _line_uses_date_as_source_evidence(line: str, mention: str) -> bool:
    escaped = re.escape(str(mention or ""))
    if not escaped:
        return False
    text = str(line or "")
    source_word = r"(?:source|sources|sourced|data source|citation|cited|evidence|dataset|primary source)"
    date_action = r"(?:published|dated|released|filed|announced)"
    patterns = (
        rf"\b{source_word}\b[^.\n;]{{0,80}}\b{date_action}\b[^.\n;]{{0,40}}\b{escaped}\b",
        rf"\b{date_action}\b[^.\n;]{{0,40}}\b{escaped}\b",
    )
    return any(re.search(pattern, text, re.I) for pattern in patterns)


def _is_constraint_planning_file(rel_path: str) -> bool:
    path = Path(rel_path)
    lowered_parts = [part.lower() for part in path.parts]
    stem = path.stem.lower()
    return bool(
        any(part in _CONSTRAINT_PLANNING_DIRS for part in lowered_parts)
        or any(token in stem for token in _CONSTRAINT_INTERNAL_NAMES)
        or any(token in stem for token in _CONSTRAINT_SUPPORT_NAMES)
    )


def _is_constraint_propagation_file(rel_path: str) -> bool:
    path = Path(rel_path)
    lowered_parts = [part.lower() for part in path.parts]
    stem = path.stem.lower()
    support_dir = any(part in _CONSTRAINT_SUPPORT_DIRS for part in lowered_parts[:-1])
    return bool(
        any(part in _CONSTRAINT_PLANNING_DIRS for part in lowered_parts)
        or any(token in stem for token in {"prompt", "subagent", "delegation"})
        or (support_dir and any(token in stem for token in {"spec", "plan"}))
    )


def _softening_lines(text: str, *, rel_path: str) -> list[str]:
    lines: list[str] = []
    if not _is_constraint_planning_file(rel_path):
        return lines
    for line in text.splitlines():
        if not _SOFTENING_RE.search(line):
            continue
        if _CONSTRAINT_CONTEXT_RE.search(line):
            lines.append(line.strip())
    return lines


def _source_date_constraint_lines(ledger: dict[str, object]) -> list[str]:
    constraints = ledger.get("constraints")
    if not isinstance(constraints, dict):
        return []
    lines: list[str] = []
    date_values = constraints.get("date")
    if isinstance(date_values, list):
        lines.extend(str(item or "").strip() for item in date_values if str(item or "").strip())
    source_values = constraints.get("source")
    if isinstance(source_values, list):
        lines.extend(
            str(item or "").strip()
            for item in source_values
            if str(item or "").strip() and _line_requires_source_date_planning_preservation(str(item or ""))
        )
    return list(dict.fromkeys(lines))


def _planning_file_preserves_source_date_constraints(text: str, constraint_lines: list[str]) -> bool:
    normalized_text = re.sub(r"\s+", " ", str(text or "")).strip().lower()
    if not normalized_text:
        return False
    if "constraint-ledger" in normalized_text or "constraint ledger" in normalized_text:
        return True
    for line in constraint_lines:
        normalized_line = re.sub(r"\s+", " ", str(line or "")).strip().lower()
        if len(normalized_line) >= 20 and normalized_line in normalized_text:
            return True
        mentions = [mention.lower() for _year, _month, mention in _month_year_mentions(line)]
        if mentions and all(mention in normalized_text for mention in mentions):
            if "source" not in normalized_line or "source" in normalized_text:
                return True
    return False


def _constraint_compliance(workspace_dir: Path, ledger: dict[str, object] | None, *, output_text: str = "") -> dict[str, object]:
    if not ledger:
        return {"status": "not_available", "issues": []}
    issues: list[dict[str, str]] = []
    warnings: list[dict[str, str]] = []
    latest_limit = _extract_latest_date_limit(ledger)
    strict = bool(ledger.get("do_not_widen_or_soften"))
    strict_source_date_lines = _source_date_constraint_lines(ledger) if strict else []
    scanned_sources: list[str] = []
    scan_entries, scan_warnings = _constraint_scan_texts(workspace_dir, output_text)
    warnings.extend(scan_warnings)
    for rel, text in scan_entries:
        scanned_sources.append(rel)
        if (
            strict_source_date_lines
            and _is_constraint_propagation_file(rel)
            and not _planning_file_preserves_source_date_constraints(text, strict_source_date_lines)
        ):
            issues.append(
                {
                    "path": rel,
                    "reason": "strict source/date constraints not referenced in planning file",
                    "text": "Planning/spec/delegation file should reference constraint-ledger or restate strict source/date constraints.",
                }
            )
        if latest_limit:
            for year, month, mention, line in _month_year_mentions_with_context(text):
                out_of_window = (
                    (year, month) > latest_limit
                    and not _is_rejected_or_out_of_scope_line(line)
                    and not _is_access_timestamp_only_line(line)
                )
                if out_of_window and _line_uses_date_as_source_evidence(line, mention):
                    issues.append(
                        {
                            "path": rel,
                            "reason": "date/source window widened past ledger limit",
                            "text": _redact_text(line.strip() or mention),
                        }
                    )
                elif (
                    strict
                    and out_of_window
                    and _is_constraint_propagation_file(rel)
                    and _CONSTRAINT_CONTEXT_RE.search(line)
                ):
                    issues.append(
                        {
                            "path": rel,
                            "reason": "planning source/date window widened past ledger limit",
                            "text": _redact_text(line.strip() or mention),
                        }
                    )
        if strict:
            for line in _softening_lines(text, rel_path=rel):
                issues.append(
                    {
                        "path": rel,
                        "reason": "strict constraint softened in workspace file",
                        "text": _redact_text(line, max_chars=240),
                    }
                )
        if len(issues) >= 100:
            break
    return {
        "status": "fail" if issues else "warn" if warnings else "pass",
        "issues": [*issues, *warnings][:100],
        "scanned_sources": scanned_sources[:100],
    }


def _required_artifact_types(ledger: dict[str, object] | None) -> list[str]:
    if not ledger:
        return []
    outputs = ledger.get("outputs")
    if not isinstance(outputs, dict):
        return []
    required: list[str] = []
    raw_formats = outputs.get("format_expectations")
    if isinstance(raw_formats, list):
        required.extend(str(item or "").strip().lower().lstrip(".") for item in raw_formats)
    raw_lines = outputs.get("required")
    if isinstance(raw_lines, list):
        required.extend(_required_output_formats([str(item or "") for item in raw_lines]))
    forbidden: list[str] = []
    raw_forbidden_formats = outputs.get("forbidden_format_expectations")
    if isinstance(raw_forbidden_formats, list):
        forbidden.extend(str(item or "").strip().lower().lstrip(".") for item in raw_forbidden_formats)
    raw_forbidden_lines = outputs.get("forbidden")
    if isinstance(raw_forbidden_lines, list):
        forbidden.extend(_output_formats([str(item or "") for item in raw_forbidden_lines]))
    forbidden_aliases: set[str] = set()
    for item in forbidden:
        forbidden_aliases.update(_format_aliases(item))
    return [item for item in dict.fromkeys(required) if item and not (_format_aliases(item) & forbidden_aliases)]


def _has_required_deliverable_intent(ledger: dict[str, object] | None) -> bool:
    if not ledger:
        return False
    outputs = ledger.get("outputs")
    if not isinstance(outputs, dict):
        return False
    raw_lines = outputs.get("required")
    if isinstance(raw_lines, list):
        required_lines = [str(item or "").strip() for item in raw_lines if str(item or "").strip()]
        if any(not _line_is_final_answer_only_context(line) for line in required_lines):
            return True
    raw_formats = outputs.get("format_expectations")
    return bool(isinstance(raw_formats, list) and any(str(item or "").strip() for item in raw_formats))


def _format_aliases(format_name: str) -> set[str]:
    normalized = str(format_name or "").strip().lower().lstrip(".")
    return _OUTPUT_FORMAT_ALIASES.get(normalized, {normalized})


def _delivered_artifact_types(artifacts: dict[str, object]) -> list[str]:
    items = artifacts.get("items") if isinstance(artifacts, dict) else []
    delivered: list[str] = []
    if not isinstance(items, list):
        return delivered
    for item in items:
        if not isinstance(item, dict):
            continue
        suffix = str(item.get("suffix") or "").strip().lower().lstrip(".")
        if suffix:
            delivered.append(suffix)
    return sorted(dict.fromkeys(delivered))


def _is_support_artifact_path(path: str) -> bool:
    parts = Path(str(path or "")).parts
    return bool(parts and parts[0].lower() in _SUPPORT_ARTIFACT_DIRS)


def _is_raw_support_capture_path(path: str) -> bool:
    parts = [part.lower() for part in Path(str(path or "")).parts]
    return bool(parts and parts[0] in _SUPPORT_ARTIFACT_DIRS and any(part in _RAW_SUPPORT_CAPTURE_DIRS for part in parts[1:]))


def _artifact_path_summary(artifacts: dict[str, object]) -> dict[str, object]:
    items = artifacts.get("items") if isinstance(artifacts, dict) else []
    if not isinstance(items, list):
        items = []
    support_paths: list[str] = []
    deliverable_paths: list[str] = []
    for item in items:
        if not isinstance(item, dict):
            continue
        rel_path = str(item.get("path") or "")
        if _is_support_artifact_path(rel_path):
            support_paths.append(rel_path)
        else:
            deliverable_paths.append(rel_path)
    return {
        "support_artifact_count": len(support_paths),
        "deliverable_artifact_count": len(deliverable_paths),
        "support_artifact_paths": support_paths[:50],
        "deliverable_artifact_paths": deliverable_paths[:50],
        "notes_only": bool(items) and not deliverable_paths and bool(support_paths),
    }


def _clean_seed_term(value: str) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip(" \t\n\r:;,.()[]{}")
    if re.fullmatch(r"<[^<>]{1,80}>", text):
        return ""
    text = re.sub(
        r"(?i)\b(?:in|with|for)\s+(?:the\s+)?(?:final\s+)?(?:report|output|artifact|deliverable)s?.*$",
        "",
        text,
    ).strip(" \t\n\r:;,.()[]{}")
    words = text.split()
    if words and words[0].lower().strip(":;,.") in _SEED_DESCRIPTOR_WORDS:
        text = " ".join(words[1:]).strip(" \t\n\r:;,.()[]{}")
    return text


def _seed_terms_from_line(line: str) -> list[str]:
    text = re.sub(r"\s+", " ", str(line or "")).strip()
    if not text:
        return []
    if re.match(r'^\s*"[A-Za-z_][A-Za-z0-9_.-]*"\s*:', text):
        return []
    candidates: list[str] = []
    quoted = re.findall(r"(?<!\w)[\"']([^\"']{2,80})[\"'](?!\w)", text)
    candidates.extend(quoted)

    body = ""
    seed_match = re.search(r"(?i)\bseed\b\s*(.*)$", text)
    if seed_match:
        body = seed_match.group(1)
    else:
        include_match = re.search(r"(?i)\b(?:must\s+)?(?:include|cover)\b\s*(.*)$", text)
        if include_match:
            body = include_match.group(1)
    if body:
        body = re.split(r"(?i)\b(?:deliver|create|output|artifact|artifacts|report|reports|source|sources)\b", body, maxsplit=1)[0]
        body = body.split(".", 1)[0]
        for part in re.split(r"\s*(?:,|;|\band\b)\s*", body):
            cleaned = _clean_seed_term(part)
            if 2 <= len(cleaned) <= 80 and len(cleaned.split()) <= 8:
                candidates.append(cleaned)

    if not candidates and ("/" in text or "\\" in text):
        if re.search(r"\b(?:seed|input|uploaded|attached|source)\s+file\b", text, re.I):
            for token in re.findall(r"[^\s,;:]+[/\\][^\s,;:]+", text):
                path = Path(token.strip(" \t\n\r:;,.()[]{}\"'"))
                if path.name:
                    candidates.append(path.name)
                    if path.stem and path.stem != path.name:
                        candidates.append(path.stem)

    terms: list[str] = []
    for candidate in candidates:
        cleaned = _clean_seed_term(candidate)
        if not cleaned:
            continue
        lower = cleaned.lower()
        if lower in _SEED_DESCRIPTOR_WORDS or lower in {"seed", "include", "cover"}:
            continue
        if lower == "confirmation" or lower.startswith("confirmation in "):
            continue
        terms.append(cleaned)
    return list(dict.fromkeys(terms))


def _seed_terms_from_ledger(ledger: dict[str, object] | None) -> list[str]:
    if not ledger:
        return []
    raw = ledger.get("seed_entities_or_files")
    if not isinstance(raw, list):
        return []
    terms: list[str] = []
    for item in raw:
        terms.extend(_seed_terms_from_line(str(item or "")))
    return list(dict.fromkeys(terms))


def _term_in_text(term: str, text: str) -> bool:
    normalized_term = " ".join(str(term or "").lower().split())
    normalized_text = " ".join(str(text or "").lower().split())
    return bool(normalized_term and normalized_term in normalized_text)


def _seed_entity_coverage(workspace_dir: Path, ledger: dict[str, object] | None, output_text: str) -> dict[str, object]:
    terms = _seed_terms_from_ledger(ledger)
    if not terms:
        return {
            "status": "not_applicable",
            "total": 0,
            "mentioned_count": 0,
            "missing": [],
            "scanned_file_count": 0,
        }
    fragments = [str(output_text or "")]
    scanned_paths: list[str] = []
    uploads_dir = workspace_dir / "uploads"
    if uploads_dir.is_dir():
        upload_fragments: list[str] = []
        for upload_path in sorted(uploads_dir.rglob("*"))[:100]:
            if not upload_path.is_file():
                continue
            try:
                rel = upload_path.relative_to(workspace_dir).as_posix()
            except ValueError:
                rel = upload_path.name
            upload_fragments.extend([rel, upload_path.name, upload_path.stem])
        if upload_fragments:
            fragments.append("\n".join(dict.fromkeys(upload_fragments)))
    for path in _iter_scan_files(workspace_dir):
        try:
            fragments.append(path.read_text(encoding="utf-8", errors="ignore"))
            scanned_paths.append(_relative_path(path, workspace_dir))
        except OSError:
            continue
    combined = "\n".join(fragments)
    mentioned = [term for term in terms if _term_in_text(term, combined)]
    missing = [term for term in terms if term not in mentioned]
    return {
        "status": "pass" if not missing else "warn",
        "total": len(terms),
        "mentioned_count": len(mentioned),
        "missing": [_redact_text(term) for term in missing[:50]],
        "scanned_file_count": len(scanned_paths),
        "scanned_paths": scanned_paths[:50],
    }


def _completion_compliance(
    workspace_dir: Path,
    artifacts: dict[str, object],
    ledger: dict[str, object] | None,
    *,
    has_final_report: bool,
    output_text: str,
) -> dict[str, object]:
    required = _required_artifact_types(ledger)
    deliverable_intent = _has_required_deliverable_intent(ledger)
    delivered = _delivered_artifact_types(artifacts)
    delivered_set = set(delivered)
    missing = [fmt for fmt in required if not (_format_aliases(fmt) & delivered_set)]
    path_summary = _artifact_path_summary(artifacts)
    seed_coverage = _seed_entity_coverage(workspace_dir, ledger, output_text)
    notes_only = bool(path_summary["notes_only"])
    issues: list[dict[str, object]] = []
    if not has_final_report:
        issues.append({"reason": "final report marker missing"})
    if missing:
        issues.append(
            {
                "reason": "required artifact types missing",
                "missing_required_artifact_types": missing,
            }
        )
    if notes_only and (required or deliverable_intent or not has_final_report):
        issues.append({"reason": "only support notes/planning artifacts were found"})
    if seed_coverage.get("status") == "warn":
        issues.append(
            {
                "reason": "seed entities/files were not all found in scanned outputs",
                "missing": seed_coverage.get("missing") or [],
            }
        )

    fail_reasons = {
        "final report marker missing",
        "required artifact types missing",
        "only support notes/planning artifacts were found",
    }
    status = "pass"
    if any(str(issue.get("reason") or "") in fail_reasons for issue in issues):
        status = "fail"
    elif issues:
        status = "warn"
    return {
        "status": status,
        "required_artifact_types": required,
        "required_deliverable_intent": deliverable_intent,
        "delivered_artifact_types": delivered,
        "missing_required_artifact_types": missing,
        "notes_only": notes_only,
        **path_summary,
        "seed_entity_coverage": seed_coverage,
        "issues": issues[:100],
    }


def _invalid_professional_artifacts(artifacts: dict[str, object]) -> list[dict[str, object]]:
    items = artifacts.get("items") if isinstance(artifacts, dict) else []
    if not isinstance(items, list):
        return []
    invalid: list[dict[str, object]] = []
    for item in items:
        if not isinstance(item, dict):
            continue
        suffix = str(item.get("suffix") or "").lower()
        if suffix not in PROFESSIONAL_ARTIFACT_EXTENSIONS:
            continue
        validation = item.get("document_validation")
        if not isinstance(validation, dict):
            continue
        if validation.get("valid") is False:
            invalid.append(
                {
                    "path": _redact_text(item.get("path") or ""),
                    "suffix": suffix,
                    "reason": _redact_text(validation.get("reason") or validation.get("library_check") or "invalid document"),
                }
            )
    return invalid


def summarize_run_evidence_result(evidence: dict[str, object]) -> dict[str, object]:
    """Summarize generic evidence checks for user-visible status/failure handling."""
    failure_reasons: list[dict[str, object]] = []
    warning_reasons: list[dict[str, object]] = []

    failure_classification = evidence.get("failure_classification")
    if isinstance(failure_classification, dict):
        failure_class = str(failure_classification.get("failure_class") or "").strip()
        if failure_class and failure_class not in {"none", "unknown"}:
            failure_reasons.append(
                {
                    "reason": "worker failure classified",
                    "failure_class": failure_class,
                    "retryable": bool(failure_classification.get("retryable")),
                    "user_message": _redact_text(failure_classification.get("user_message") or ""),
                }
            )

    completion = evidence.get("completion_compliance")
    final_output = evidence.get("final_output")
    if isinstance(final_output, dict) and not final_output.get("has_final_report"):
        failure_reasons.append({"reason": "final report marker missing"})

    if isinstance(completion, dict):
        completion_status = str(completion.get("status") or "").lower()
        issues = completion.get("issues") if isinstance(completion.get("issues"), list) else []
        if completion_status == "fail":
            failure_reasons.append({"reason": "completion compliance failed", "issues": issues[:10]})
        elif completion_status == "warn":
            warning_reasons.append({"reason": "completion compliance warning", "issues": issues[:10]})

    coverage = evidence.get("coverage_compliance")
    if isinstance(coverage, dict):
        coverage_status = str(coverage.get("status") or "").lower()
        issues = coverage.get("issues") if isinstance(coverage.get("issues"), list) else []
        if coverage_status == "fail":
            failure_reasons.append({"reason": "coverage compliance failed", "issues": issues[:10]})
        elif coverage_status == "warn":
            warning_reasons.append({"reason": "coverage compliance warning", "issues": issues[:10]})

    constraint = evidence.get("constraint_compliance")
    if isinstance(constraint, dict):
        constraint_status = str(constraint.get("status") or "").lower()
        issues = constraint.get("issues") if isinstance(constraint.get("issues"), list) else []
        if constraint_status == "fail":
            failure_reasons.append({"reason": "constraint compliance failed", "issues": issues[:10]})
        elif constraint_status == "warn":
            warning_reasons.append({"reason": "constraint compliance warning", "issues": issues[:10]})

    invalid_artifacts = _invalid_professional_artifacts(
        evidence.get("artifacts") if isinstance(evidence.get("artifacts"), dict) else {}
    )
    if invalid_artifacts:
        failure_reasons.append(
            {
                "reason": "professional artifact validation failed",
                "artifacts": invalid_artifacts[:20],
            }
        )

    hygiene = evidence.get("content_hygiene")
    if isinstance(hygiene, dict) and str(hygiene.get("status") or "").lower() in {"warn", "fail"}:
        warning_reasons.append(
            {
                "reason": "content hygiene warning",
                "failure_count": hygiene.get("failure_count"),
                "items": (hygiene.get("items") if isinstance(hygiene.get("items"), list) else [])[:10],
            }
        )

    status = "pass"
    if failure_reasons:
        status = "fail"
    elif warning_reasons:
        status = "warn"
    return {
        "status": status,
        "failure_reasons": failure_reasons[:20],
        "warning_reasons": warning_reasons[:20],
    }


def _failure_classification_summary(
    *,
    stdout_text: str,
    stderr_text: str,
    runtime_name: str,
    exit_code: int | None,
    error_text: str,
    has_final_report: bool,
) -> dict[str, object]:
    if (
        exit_code == 0
        and not str(error_text or "").strip()
        and (has_final_report or not has_structured_failure_evidence(stdout_text or "", stderr_text or ""))
    ):
        return {"status": "not_applicable"}
    classification = classify_cli_failure(
        stdout=stdout_text or "",
        stderr=stderr_text or error_text or "",
        runtime_name=runtime_name,
        exit_code=exit_code,
    )
    return {
        "status": "classified",
        "failure_class": classification.failure_class,
        "retryable": classification.retryable,
        "user_message": _redact_text(classification.user_message),
        "recommended_recovery": _redact_text(classification.recommended_recovery),
        "diagnostic_summary": _redact_text(classification.diagnostic_summary, max_chars=1200),
    }


def build_run_evidence(
    *,
    worker: dict[str, object],
    run_id: str,
    runtime_name: str,
    model: str,
    command: list[str],
    env: dict[str, str] | None,
    workspace_dir: Path | str,
    stdout_text: str,
    stderr_text: str,
    output_text: str,
    error_text: str,
    exit_code: int | None,
    timeout_seconds: float | None,
    stop_reason: str,
    constraint_ledger: dict[str, object] | None,
    started_at: float | None = None,
    ended_at: float | None = None,
    transcript_paths: dict[str, str] | None = None,
) -> dict[str, object]:
    workspace = Path(workspace_dir)
    artifacts = _artifact_inventory(workspace, started_at=started_at, ended_at=ended_at)
    normalized_stop_reason = str(stop_reason or "").strip()
    if normalized_stop_reason == "timeout":
        exit_source = "timeout"
    elif exit_code is not None:
        exit_source = "process"
    else:
        exit_source = normalized_stop_reason or "unknown"
    has_final_report = bool(
        _FINAL_REPORT_RE.search(output_text)
        or _stdout_agent_has_final_report(stdout_text)
    )
    final_output = {
        "has_final_report": has_final_report,
        "output_chars": len(output_text or ""),
        "error_present": bool(str(error_text or "").strip()),
        "status": "ok" if exit_code == 0 and not str(error_text or "").strip() else "failed",
    }
    effective_effort = _effective_effort(worker, command, env)
    evidence: dict[str, object] = {
        "schema": "glasshive.run.evidence.v1",
        "run_id": run_id,
        "created_at": _now_epoch(),
        "started_at": started_at,
        "ended_at": ended_at if ended_at is not None else _now_epoch(),
        "worker": {
            "worker_id": str(worker.get("worker_id") or ""),
            "profile": str(worker.get("profile") or ""),
            "execution_mode": str(worker.get("execution_mode") or ""),
            "runtime": str(worker.get("runtime") or ""),
            "backend": _derived_legacy_backend(worker),
        },
        "runtime": runtime_name,
        "model": model,
        "effort": effective_effort,
        "effort_projection": _effort_projection(worker, effective_effort),
        "command": {
            "argv_redacted": [_redact_command_arg(part) for part in command],
            "display_redacted": " ".join(_redact_command_arg(part) for part in command),
        },
        "env_keys": _safe_env_keys(env),
        "redacted_env_key_count": len((env or {}).keys()) - len(_safe_env_keys(env)),
        "timeout": {
            "seconds": timeout_seconds,
            "exit_source": exit_source,
            "stop_reason": normalized_stop_reason,
        },
        "transcript": {
            "paths": transcript_paths or {},
            "metadata": _transcript_metadata(transcript_paths, workspace_dir=workspace),
            "stdout_tail": _redact_text(stdout_text, max_chars=4000),
            "stderr_tail": _redact_text(stderr_text, max_chars=4000),
        },
        "artifacts": artifacts,
        "visual_render_evidence": _visual_render_summary(artifacts),
        "content_hygiene": check_content_hygiene(_text_artifact_payload(workspace)),
        "constraint_compliance": _constraint_compliance(workspace, constraint_ledger, output_text=output_text),
        "coverage_compliance": _coverage_compliance(
            workspace,
            artifacts,
            constraint_ledger,
            output_text=output_text,
        ),
        "completion_compliance": _completion_compliance(
            workspace,
            artifacts,
            constraint_ledger,
            has_final_report=has_final_report,
            output_text=output_text,
        ),
        "final_output": final_output,
        "failure_classification": _failure_classification_summary(
            stdout_text=stdout_text,
            stderr_text=stderr_text,
            runtime_name=runtime_name,
            exit_code=exit_code,
            error_text=error_text,
            has_final_report=has_final_report,
        ),
        "exit_code": exit_code,
    }
    evidence["evidence_result"] = summarize_run_evidence_result(evidence)
    return evidence


def write_run_evidence(workspace_dir: Path | str, evidence: dict[str, object], run_id: str) -> Path:
    run_dir = Path(workspace_dir) / RUN_EVIDENCE_DIR_NAME
    latest_path = run_dir / "evidence.json"
    per_run_path = run_dir / "runs" / run_id / "evidence.json"
    latest_path.parent.mkdir(parents=True, exist_ok=True)
    per_run_path.parent.mkdir(parents=True, exist_ok=True)
    body = json.dumps(evidence, indent=2, sort_keys=True)
    latest_path.write_text(body + "\n")
    per_run_path.write_text(body + "\n")
    return latest_path
